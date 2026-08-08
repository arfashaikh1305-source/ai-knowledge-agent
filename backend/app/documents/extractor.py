from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
import markdown


def extract_text(file_path: str) -> str:
    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        return extract_pdf(file_path)

    elif extension == ".docx":
        return extract_docx(file_path)

    elif extension == ".pptx":
        return extract_pptx(file_path)

    elif extension == ".xlsx":
        return extract_xlsx(file_path)

    elif extension == ".txt":
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()

    elif extension == ".md":
        return extract_markdown(file_path)

    return ""


def extract_pdf(file_path: str):
    reader = PdfReader(file_path)

    text = ""

    for page in reader.pages:
        page_text = page.extract_text()

        if page_text:
            text += page_text + "\n"

    return text


def extract_docx(file_path: str):
    doc = DocxDocument(file_path)

    return "\n".join(
        paragraph.text
        for paragraph in doc.paragraphs
    )


def extract_pptx(file_path: str):
    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def extract_xlsx(file_path: str):
    workbook = load_workbook(file_path)

    text = ""

    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            values = [
                str(cell)
                for cell in row
                if cell is not None
            ]

            if values:
                text += " ".join(values) + "\n"

    return text


def extract_markdown(file_path: str):
    with open(file_path, "r", encoding="utf-8") as file:
        md_text = file.read()

    # Convert Markdown to plain text (HTML output)
    html = markdown.markdown(md_text)

    return html